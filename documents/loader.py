"""
NeuraDocs - Universal Document Loader
=====================================
Multi-format file parser for PDF, DOCX, TXT, MD, CSV, XLSX, PPTX, HTML,
with automatic OCR fallback for scanned PDFs.
"""

from __future__ import annotations

import io
import re
from dataclasses import dataclass
from typing import Dict, List, Optional, Tuple


@dataclass
class ExtractedPage:
    page_number: int
    text: str
    has_tables: bool = False
    section_name: Optional[str] = None


@dataclass
class ExtractedDocument:
    filename: str
    file_type: str
    pages: List[ExtractedPage]
    total_characters: int
    metadata: Dict[str, any]

    @property
    def full_text(self) -> str:
        return "\n\n".join(p.text for p in self.pages if p.text.strip())


def extract_document(file_bytes: bytes, filename: str) -> ExtractedDocument:
    """Universal parser that inspects extension and returns structured ExtractedDocument."""
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""

    if ext == "pdf":
        return _extract_pdf(file_bytes, filename)
    elif ext == "docx":
        return _extract_docx(file_bytes, filename)
    elif ext == "pptx":
        return _extract_pptx(file_bytes, filename)
    elif ext in ("csv", "tsv"):
        return _extract_tabular(file_bytes, filename, delimiter="," if ext == "csv" else "\t")
    elif ext in ("xlsx", "xls"):
        return _extract_excel(file_bytes, filename)
    elif ext in ("txt", "md", "markdown", "py", "json", "yaml", "yml"):
        return _extract_plain_text(file_bytes, filename, ext)
    elif ext in ("html", "htm"):
        return _extract_html(file_bytes, filename)
    else:
        # Fallback to UTF-8 decoding
        try:
            return _extract_plain_text(file_bytes, filename, ext)
        except Exception:
            raise ValueError(f"Unsupported file format '.{ext}'. Supported: PDF, DOCX, PPTX, TXT, MD, CSV, XLSX, HTML")


def _extract_pdf(file_bytes: bytes, filename: str) -> ExtractedDocument:
    pages: List[ExtractedPage] = []
    
    # 1. Try pdfplumber first for superior text layout and table extraction
    try:
        import pdfplumber

        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for idx, page in enumerate(pdf.pages, start=1):
                page_text = ""
                has_tables = False
                
                # Extract structured tables if present
                tables = page.extract_tables()
                if tables:
                    has_tables = True
                    for table in tables:
                        table_str = _format_table_as_markdown(table)
                        if table_str:
                            page_text += f"\n\n[TABLE]\n{table_str}\n[/TABLE]\n\n"

                # Extract text
                extracted = page.extract_text(layout=True) or ""
                if extracted.strip():
                    page_text += extracted
                
                pages.append(ExtractedPage(page_number=idx, text=page_text.strip(), has_tables=has_tables))
    except Exception:
        # 2. Fallback to pypdf if pdfplumber fails
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(file_bytes))
        pages = []
        for idx, page in enumerate(reader.pages, start=1):
            try:
                t = page.extract_text() or ""
            except Exception:
                t = ""
            pages.append(ExtractedPage(page_number=idx, text=t.strip(), has_tables=False))

    # 3. Check for scanned / empty pages and attempt OCR if pytesseract is available
    total_len = sum(len(p.text) for p in pages)
    if total_len < 50 and len(pages) > 0:
        try:
            import pdf2image
            import pytesseract
            from PIL import Image

            images = pdf2image.convert_from_bytes(file_bytes, dpi=200)
            ocr_pages = []
            for idx, img in enumerate(images, start=1):
                ocr_text = pytesseract.image_to_string(img)
                ocr_pages.append(ExtractedPage(page_number=idx, text=ocr_text.strip(), has_tables=False))
            if sum(len(p.text) for p in ocr_pages) > total_len:
                pages = ocr_pages
        except Exception:
            # OCR optional fallback silently passes if tesseract binary is not installed
            pass

    total_chars = sum(len(p.text) for p in pages)
    return ExtractedDocument(
        filename=filename,
        file_type="pdf",
        pages=pages,
        total_characters=total_chars,
        metadata={"num_pages": len(pages)},
    )


def _extract_pptx(file_bytes: bytes, filename: str) -> ExtractedDocument:
    """Extract text from PowerPoint presentations slide-by-slide."""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation(io.BytesIO(file_bytes))
    pages: List[ExtractedPage] = []

    for idx, slide in enumerate(prs.slides, start=1):
        parts = []
        # Slide title
        if slide.shapes.title and slide.shapes.title.text.strip():
            parts.append(f"### {slide.shapes.title.text.strip()}")
        # All text boxes
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text and text != (slide.shapes.title.text.strip() if slide.shapes.title else ""):
                        parts.append(text)
        # Speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                parts.append(f"[Notes: {notes_text}]")

        slide_text = "\n".join(parts).strip()
        if slide_text:
            pages.append(ExtractedPage(page_number=idx, text=slide_text, has_tables=False))

    total_chars = sum(len(p.text) for p in pages)
    return ExtractedDocument(
        filename=filename,
        file_type="pptx",
        pages=pages,
        total_characters=total_chars,
        metadata={"slides": len(prs.slides)},
    )


def _extract_docx(file_bytes: bytes, filename: str) -> ExtractedDocument:
    import docx

    doc = docx.Document(io.BytesIO(file_bytes))
    parts = []
    has_tables = False

    # Extract headings and paragraphs with structural cues
    for p in doc.paragraphs:
        if p.text.strip():
            if p.style.name.startswith("Heading"):
                parts.append(f"\n### {p.text.strip()}\n")
            else:
                parts.append(p.text.strip())

    # Extract tables in DOCX
    for table in doc.tables:
        has_tables = True
        table_data = []
        for row in table.rows:
            table_data.append([cell.text.strip() for cell in row.cells])
        md_table = _format_table_as_markdown(table_data)
        if md_table:
            parts.append(f"\n[TABLE]\n{md_table}\n[/TABLE]\n")

    full_text = "\n\n".join(parts).strip()
    return ExtractedDocument(
        filename=filename,
        file_type="docx",
        pages=[ExtractedPage(page_number=1, text=full_text, has_tables=has_tables)],
        total_characters=len(full_text),
        metadata={"paragraphs": len(doc.paragraphs), "tables": len(doc.tables)},
    )


def _extract_tabular(file_bytes: bytes, filename: str, delimiter: str = ",") -> ExtractedDocument:
    import pandas as pd

    df = pd.read_csv(io.BytesIO(file_bytes), sep=delimiter)
    # Convert to Markdown table summary
    md_repr = df.to_markdown(index=False)
    summary = f"Dataset: {filename} ({len(df)} rows, {len(df.columns)} columns)\nColumns: {', '.join(df.columns)}\n\n{md_repr}"

    return ExtractedDocument(
        filename=filename,
        file_type="csv",
        pages=[ExtractedPage(page_number=1, text=summary, has_tables=True)],
        total_characters=len(summary),
        metadata={"rows": len(df), "cols": len(df.columns)},
    )


def _extract_excel(file_bytes: bytes, filename: str) -> ExtractedDocument:
    import pandas as pd

    excel_file = pd.ExcelFile(io.BytesIO(file_bytes))
    pages = []

    for idx, sheet_name in enumerate(excel_file.sheet_names, start=1):
        df = pd.read_excel(excel_file, sheet_name=sheet_name)
        md_repr = df.to_markdown(index=False) if len(df) <= 500 else df.head(500).to_markdown(index=False)
        sheet_text = (
            f"### Sheet: {sheet_name}\n"
            f"Dimensions: {len(df)} rows, {len(df.columns)} columns\n"
            f"Columns: {', '.join(str(c) for c in df.columns)}\n\n"
            f"{md_repr}"
        )
        pages.append(
            ExtractedPage(
                page_number=idx,
                text=sheet_text,
                has_tables=True,
                section_name=sheet_name,
            )
        )

    total_chars = sum(len(p.text) for p in pages)
    return ExtractedDocument(
        filename=filename,
        file_type="excel",
        pages=pages,
        total_characters=total_chars,
        metadata={"sheets": excel_file.sheet_names},
    )


def _extract_plain_text(file_bytes: bytes, filename: str, ext: str) -> ExtractedDocument:
    text = file_bytes.decode("utf-8", errors="ignore")
    return ExtractedDocument(
        filename=filename,
        file_type=ext or "txt",
        pages=[ExtractedPage(page_number=1, text=text.strip(), has_tables=False)],
        total_characters=len(text),
        metadata={},
    )


def _extract_html(file_bytes: bytes, filename: str) -> ExtractedDocument:
    import re

    raw_html = file_bytes.decode("utf-8", errors="ignore")
    # Clean scripts, styles
    cleaned = re.sub(r"<(script|style).*?</\1>", "", raw_html, flags=re.DOTALL | re.IGNORECASE)
    # Strip HTML tags
    text = re.sub(r"<[^>]+>", " ", cleaned)
    text = re.sub(r"\s+", " ", text).strip()

    return ExtractedDocument(
        filename=filename,
        file_type="html",
        pages=[ExtractedPage(page_number=1, text=text, has_tables=False)],
        total_characters=len(text),
        metadata={},
    )


def _format_table_as_markdown(table: List[List[Optional[str]]]) -> str:
    if not table or len(table) < 2:
        return ""

    cleaned_rows = []
    for row in table:
        cleaned_row = [re.sub(r"\s+", " ", str(cell or "").strip()) for cell in row]
        if any(cleaned_row):
            cleaned_rows.append(cleaned_row)

    if not cleaned_rows:
        return ""

    headers = cleaned_rows[0]
    num_cols = len(headers)
    separator = ["---"] * num_cols

    md_lines = [
        "| " + " | ".join(headers) + " |",
        "| " + " | ".join(separator) + " |",
    ]

    for row in cleaned_rows[1:]:
        # Pad or truncate row to match headers count
        padded_row = (row + [""] * num_cols)[:num_cols]
        md_lines.append("| " + " | ".join(padded_row) + " |")

    return "\n".join(md_lines)
